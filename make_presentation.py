"""
Generate WSCprep introductory presentation as a PowerPoint file.
Run: python make_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0F, 0x2E, 0x5C)   # deep navy  – title bg
TEAL       = RGBColor(0x00, 0x7A, 0x8A)   # teal       – accent / headings
LIGHT_TEAL = RGBColor(0xE0, 0xF4, 0xF7)   # light teal – body bg
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY  = RGBColor(0x2D, 0x2D, 0x2D)
MID_GRAY   = RGBColor(0x55, 0x65, 0x70)
LIGHT_GRAY = RGBColor(0xF2, 0xF5, 0xF7)
ORANGE     = RGBColor(0xE8, 0x6A, 0x17)   # highlight

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # truly blank layout

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None, line_width=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = line_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, font_size=16, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, space_before=Pt(4), italic=False):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size   = _Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def slide_header(slide, title, subtitle=None):
    """Adds a coloured header bar with title + optional subtitle."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), NAVY)
    add_text(slide, title,
             Inches(0.5), Inches(0.1), Inches(12), Inches(0.7),
             font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.78), Inches(12), Inches(0.4),
                 font_size=15, color=LIGHT_TEAL)
    # thin teal accent line under header
    add_rect(slide, 0, Inches(1.25), SLIDE_W, Pt(3), TEAL)


def body_bg(slide):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)


def feature_card(slide, x, y, w, h, icon, heading, bullets):
    """A rounded-corner-style feature card (rectangle + text)."""
    add_rect(slide, x, y, w, h, WHITE, line_rgb=TEAL, line_width=Pt(1.5))
    # icon + heading
    add_text(slide, f"{icon}  {heading}",
             x + Inches(0.15), y + Inches(0.1),
             w - Inches(0.3), Inches(0.42),
             font_size=14, bold=True, color=TEAL)
    # bullet lines
    by = y + Inches(0.52)
    bh = Inches(0.28)
    for b in bullets:
        add_text(slide, f"  • {b}",
                 x + Inches(0.15), by,
                 w - Inches(0.3), bh,
                 font_size=11, color=DARK_GRAY)
        by += bh


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)

# Full navy background
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)

# Teal accent bar (left edge)
add_rect(s, 0, 0, Inches(0.18), SLIDE_H, TEAL)

# White decorative strip at bottom
add_rect(s, 0, Inches(6.6), SLIDE_W, Inches(0.9), TEAL)

# Title
add_text(s, "WSC Hydrometric\nAnalysis App",
         Inches(0.55), Inches(1.4), Inches(9), Inches(2.5),
         font_size=48, bold=True, color=WHITE)

# Subtitle
add_text(s, "Interactive web platform for Water Survey of Canada\nstream-flow data analysis and flood frequency modelling",
         Inches(0.55), Inches(3.9), Inches(9.5), Inches(1.2),
         font_size=20, color=LIGHT_TEAL, italic=True)

# URL badge
add_rect(s, Inches(0.55), Inches(5.35), Inches(4.4), Inches(0.6), ORANGE)
add_text(s, "  🌐  ws-cprep.vercel.app",
         Inches(0.6), Inches(5.38), Inches(4.3), Inches(0.55),
         font_size=17, bold=True, color=WHITE)

# Bottom strip text
add_text(s, "Powered by ECCC HYDAT  ·  Built with Next.js + Python FastAPI",
         Inches(0.3), Inches(6.62), Inches(12), Inches(0.65),
         font_size=13, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — WHAT IS IT?
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
body_bg(s)
slide_header(s, "What Is This App?",
             "A one-stop platform for WSC hydrometric data exploration & analysis")

add_rect(s, Inches(0.4), Inches(1.5), Inches(12.5), Inches(5.5), WHITE,
         line_rgb=LIGHT_TEAL)

txb = s.shapes.add_textbox(Inches(0.7), Inches(1.65), Inches(11.9), Inches(5))
tf  = txb.text_frame
tf.word_wrap = True

add_para(tf,
    "Water Survey of Canada (WSC) operates ~2,000 active hydrometric stations across the country. "
    "Until now, extracting and analysing that data meant juggling the HYDAT SQLite database, R scripts, "
    "or bulky desktop GIS tools.",
    font_size=16, color=DARK_GRAY, space_before=Pt(0))

add_para(tf, "", font_size=8)

add_para(tf,
    "This app replaces that workflow with a modern, browser-based platform that anyone on your team "
    "can open without installing anything. It pulls live data directly from Environment and Climate "
    "Change Canada's OGC API, runs all analyses on the fly, and presents results in clear, "
    "publication-ready charts.",
    font_size=16, color=DARK_GRAY)

add_para(tf, "", font_size=8)

add_para(tf, "Key facts:", font_size=15, bold=True, color=TEAL)
for fact in [
    "✅   No software installation — runs entirely in the browser",
    "✅   Live data from ECCC's national HYDAT database (updated daily)",
    "✅   Covers all active & historical WSC stream-flow & water-level stations",
    "✅   Analyses from basic statistics to advanced flood frequency modelling",
    "✅   Export charts, tables, and printable PDF reports",
]:
    add_para(tf, fact, font_size=15, color=DARK_GRAY, space_before=Pt(5))


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — STATION DISCOVERY
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
body_bg(s)
slide_header(s, "Station Discovery",
             "Find and explore any WSC hydrometric station instantly")

# Left panel — description
add_rect(s, Inches(0.4), Inches(1.45), Inches(6.0), Inches(5.6), WHITE)

txb = s.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(5.6), Inches(5.2))
tf  = txb.text_frame; tf.word_wrap = True

add_para(tf, "Interactive Map", font_size=18, bold=True, color=TEAL, space_before=Pt(0))
add_para(tf,
    "The home screen displays every WSC station on an interactive Leaflet map. "
    "Switch between OpenStreetMap and satellite imagery with one click.",
    font_size=14, color=DARK_GRAY)

add_para(tf, "", font_size=6)
add_para(tf, "Colour-coded by metric", font_size=14, bold=True, color=TEAL)
add_para(tf,
    "Station markers can be coloured by mean annual discharge, record length, "
    "or drainage area — giving an instant spatial overview of the network.",
    font_size=14, color=DARK_GRAY)

add_para(tf, "", font_size=6)
add_para(tf, "Search & Filter", font_size=14, bold=True, color=TEAL)
for item in [
    "Search by station name or ID",
    "Filter by province or active/historical status",
    "Save favourite stations for quick access",
]:
    add_para(tf, f"  • {item}", font_size=13, color=DARK_GRAY, space_before=Pt(3))

add_para(tf, "", font_size=6)
add_para(tf, "Station Details", font_size=14, bold=True, color=TEAL)
add_para(tf,
    "Clicking any station opens its full profile: drainage area, record span, "
    "current status, and links to all analysis tabs.",
    font_size=14, color=DARK_GRAY)

# Right panel — visual placeholder
add_rect(s, Inches(6.65), Inches(1.45), Inches(6.3), Inches(5.6), NAVY)
add_text(s, "🗺️",
         Inches(8.8), Inches(2.8), Inches(2), Inches(2), font_size=60)
add_text(s,
    "Interactive map showing\n~2,000 WSC stations\ncolour-coded by metric",
    Inches(6.75), Inches(4.8), Inches(6.1), Inches(1.6),
    font_size=15, color=LIGHT_TEAL, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — HYDROGRAPH & BASEFLOW
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
body_bg(s)
slide_header(s, "Hydrograph & Baseflow Separation",
             "Visualise the full daily record and decompose flow components")

col_w = Inches(5.9)
gap   = Inches(0.5)
y0    = Inches(1.45)
h     = Inches(5.6)

# Card 1 – Hydrograph
add_rect(s, Inches(0.4), y0, col_w, h, WHITE)
txb = s.shapes.add_textbox(Inches(0.6), y0+Inches(0.15), col_w-Inches(0.3), h-Inches(0.2))
tf  = txb.text_frame; tf.word_wrap = True
add_para(tf, "📈  Daily Hydrograph", font_size=18, bold=True, color=TEAL, space_before=Pt(0))
add_para(tf,
    "Plot the complete daily discharge or water-level record for any station. "
    "Pan and zoom across decades of data with an interactive timeline.",
    font_size=14, color=DARK_GRAY)
add_para(tf, "", font_size=6)
add_para(tf, "Features:", font_size=13, bold=True, color=MID_GRAY)
for f in [
    "Date-range selector (zoom to any period)",
    "Toggle discharge ↔ water level",
    "High-flow event flags on the chart",
    "Download raw data as CSV",
    "Variable colour by flow magnitude",
]:
    add_para(tf, f"  • {f}", font_size=13, color=DARK_GRAY, space_before=Pt(3))

# Card 2 – Baseflow
add_rect(s, Inches(0.4)+col_w+gap, y0, col_w, h, WHITE)
txb = s.shapes.add_textbox(Inches(0.6)+col_w+gap, y0+Inches(0.15), col_w-Inches(0.3), h-Inches(0.2))
tf  = txb.text_frame; tf.word_wrap = True
add_para(tf, "🌊  Baseflow Separation", font_size=18, bold=True, color=TEAL, space_before=Pt(0))
add_para(tf,
    "Decompose total discharge into baseflow (groundwater) and quickflow (storm runoff) "
    "using two industry-standard digital filters — both running entirely in your browser.",
    font_size=14, color=DARK_GRAY)
add_para(tf, "", font_size=6)
add_para(tf, "Two methods:", font_size=13, bold=True, color=MID_GRAY)
add_para(tf, "  Lyne & Hollick (1979)", font_size=13, bold=True, color=DARK_GRAY, space_before=Pt(4))
add_para(tf, "  3-pass recursive filter (α = 0.925)", font_size=12, color=MID_GRAY, space_before=Pt(2))
add_para(tf, "", font_size=4)
add_para(tf, "  Eckhardt (2005)", font_size=13, bold=True, color=DARK_GRAY, space_before=Pt(4))
add_para(tf, "  Two-parameter filter (a = 0.98, BFImax = 0.80)", font_size=12, color=MID_GRAY, space_before=Pt(2))
add_para(tf, "", font_size=6)
add_para(tf, "Output:", font_size=13, bold=True, color=MID_GRAY)
add_para(tf,
    "Stacked-area chart showing total flow, baseflow, and quickflow bands. "
    "Live Baseflow Index (BFI) readout updates as you adjust filter parameters.",
    font_size=13, color=DARK_GRAY, space_before=Pt(4))

# Right label block
add_rect(s, Inches(12.1), y0, Inches(0.83), h, TEAL)
add_text(s, "HYDROGRAPH", Inches(12.1), Inches(3.4), Inches(0.83), Inches(2),
         font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — FDC + REGIME + STATS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
body_bg(s)
slide_header(s, "Flow Analysis Tabs",
             "Flow Duration Curve · Annual Regime · Statistical Summary")

CW = Inches(3.9)
CH = Inches(5.5)
Y0 = Inches(1.5)
GAP = Inches(0.2)
X = [Inches(0.35), Inches(0.35)+CW+GAP, Inches(0.35)+2*(CW+GAP)]

cards = [
    ("📉", "Flow Duration Curve",
     "Ranks every daily discharge from highest to lowest and plots "
     "exceedance probability on a log scale.",
     [
         "Q10, Q50, Q90 markers",
         "Identifies high-flow, median,\nand low-flow regimes",
         "Essential for environmental\nflow assessments",
         "Log-scale discharge axis",
     ]),
    ("📅", "Annual Regime",
     "Shows the seasonal pattern of flow throughout the year — "
     "averaged across the full station record.",
     [
         "Mean monthly discharge",
         "Q10–Q90 uncertainty band",
         "Reveals snowmelt peaks,\nsummer low-flows, etc.",
         "Useful for water allocation\nplanning",
     ]),
    ("📊", "Statistical Summary",
     "Compact table of key descriptive statistics for the complete "
     "daily record at a glance.",
     [
         "Mean, median, std deviation",
         "Min, Max, CV",
         "Record length & completeness",
         "Annual peak statistics",
     ]),
]

for i, (icon, title, desc, bullets) in enumerate(cards):
    add_rect(s, X[i], Y0, CW, CH, WHITE)
    txb = s.shapes.add_textbox(X[i]+Inches(0.15), Y0+Inches(0.1),
                                CW-Inches(0.3), CH-Inches(0.2))
    tf = txb.text_frame; tf.word_wrap = True
    add_para(tf, f"{icon}  {title}", font_size=16, bold=True, color=TEAL, space_before=Pt(0))
    add_para(tf, desc, font_size=13, color=DARK_GRAY, space_before=Pt(6))
    add_para(tf, "", font_size=5)
    for b in bullets:
        add_para(tf, f"  • {b}", font_size=12, color=DARK_GRAY, space_before=Pt(4))


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — TREND ANALYSIS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
body_bg(s)
slide_header(s, "Trend Analysis",
             "Detect statistically significant long-term changes in streamflow")

add_rect(s, Inches(0.4), Inches(1.45), Inches(12.5), Inches(5.6), WHITE)

txb = s.shapes.add_textbox(Inches(0.65), Inches(1.6), Inches(11.9), Inches(5.2))
tf  = txb.text_frame; tf.word_wrap = True

add_para(tf, "Mann-Kendall Trend Test", font_size=20, bold=True, color=TEAL, space_before=Pt(0))
add_para(tf,
    "The Mann-Kendall test is the standard non-parametric method for detecting monotonic trends "
    "in hydrological time series. It makes no assumptions about the underlying distribution of the data, "
    "making it robust even for skewed or non-normal streamflow records.",
    font_size=15, color=DARK_GRAY, space_before=Pt(8))

add_para(tf, "", font_size=5)
add_para(tf, "What the app reports:", font_size=15, bold=True, color=TEAL)

cols = [
    ("Trend Direction", "Increasing / Decreasing / No trend — determined at the selected significance level (α = 0.05 or 0.01)."),
    ("Sen's Slope", "Magnitude of the trend in m³/s per year (or m/yr for water level). A median-based robust slope estimator."),
    ("p-value", "Probability of observing the data under the null hypothesis of no trend. Values < 0.05 are statistically significant."),
    ("Annual Peaks Trend", "Separate analysis on the annual maximum series to detect changes in flood magnitude over time."),
]

for label, desc in cols:
    add_para(tf, f"  ▸  {label}:", font_size=14, bold=True, color=NAVY, space_before=Pt(8))
    add_para(tf, f"      {desc}", font_size=13, color=DARK_GRAY, space_before=Pt(2))

add_para(tf, "", font_size=5)
add_para(tf,
    "Results are displayed with an annual time series plot and a Sen's slope regression line, "
    "so the direction and rate of change are immediately visible.",
    font_size=14, color=MID_GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — FLOOD FREQUENCY ANALYSIS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
body_bg(s)
slide_header(s, "Flood Frequency Analysis (FFA)",
             "Estimate design floods for any return period using multiple distributions")

# Left – method description
add_rect(s, Inches(0.4), Inches(1.45), Inches(7.4), Inches(5.6), WHITE)
txb = s.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(7.1), Inches(5.2))
tf  = txb.text_frame; tf.word_wrap = True

add_para(tf, "Annual Maximum Series (AMS) approach", font_size=16, bold=True, color=TEAL, space_before=Pt(0))
add_para(tf,
    "Fits multiple probability distributions to the annual peak flow series "
    "and selects the best fit using information criteria (AIC/BIC).",
    font_size=14, color=DARK_GRAY)

add_para(tf, "", font_size=5)
add_para(tf, "Distributions fitted:", font_size=14, bold=True, color=TEAL)
dists = [
    ("GEV",    "Generalised Extreme Value — most widely used for floods"),
    ("Gumbel", "Special case of GEV; common in older Canadian practice"),
    ("LP3",    "Log-Pearson Type III — standard in US & Canadian guidelines"),
    ("LN2",    "2-parameter Lognormal"),
    ("Pearson 3", "Generalised gamma family"),
]
for abbr, desc in dists:
    add_para(tf, f"  • {abbr}  —  {desc}", font_size=13, color=DARK_GRAY, space_before=Pt(4))

add_para(tf, "", font_size=5)
add_para(tf, "Parameter estimation:", font_size=14, bold=True, color=TEAL)
add_para(tf, "  L-moments (default) or Maximum Likelihood (MLE)", font_size=13, color=DARK_GRAY, space_before=Pt(4))

add_para(tf, "", font_size=5)
add_para(tf, "Confidence intervals:", font_size=14, bold=True, color=TEAL)
add_para(tf, "  90% bootstrap CI on all quantile estimates", font_size=13, color=DARK_GRAY, space_before=Pt(4))

# Right – outputs
add_rect(s, Inches(8.05), Inches(1.45), Inches(4.9), Inches(5.6), NAVY)
txb = s.shapes.add_textbox(Inches(8.25), Inches(1.6), Inches(4.55), Inches(5.2))
tf  = txb.text_frame; tf.word_wrap = True

add_para(tf, "Outputs", font_size=18, bold=True, color=WHITE, space_before=Pt(0))
add_para(tf, "", font_size=5)

outputs = [
    ("📈  Frequency curve",
     "Empirical plotting positions + fitted distribution lines on a log return-period axis"),
    ("📋  Design flood table",
     "Quantiles (m³/s) for T = 2, 5, 10, 20, 50, 100, 200, 500 years"),
    ("📐  Goodness-of-fit",
     "AIC / BIC scores and K-S test p-value for each distribution"),
    ("🔄  Plotting positions",
     "Choice of Weibull, Hazen, or Cunnane"),
    ("⚙️  Estimated-year toggle",
     "Option to exclude provisional or estimated records"),
]
for icon_title, desc in outputs:
    add_para(tf, icon_title, font_size=14, bold=True, color=LIGHT_TEAL, space_before=Pt(10))
    add_para(tf, desc, font_size=12, color=WHITE, space_before=Pt(2))


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — PEAKS OVER THRESHOLD (POT)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
body_bg(s)
slide_header(s, "Peaks Over Threshold (POT) Analysis",
             "Partial duration series — more data points, better estimates for short records")

add_rect(s, Inches(0.4), Inches(1.45), Inches(12.5), Inches(5.6), WHITE)
txb = s.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(11.9), Inches(5.2))
tf  = txb.text_frame; tf.word_wrap = True

add_para(tf, "Why POT?", font_size=18, bold=True, color=TEAL, space_before=Pt(0))
add_para(tf,
    "Traditional flood frequency analysis uses one peak per year (Annual Maximum Series). "
    "POT uses every independent flow event above a user-selected threshold — often 3–5× more data "
    "points for the same record length. This is especially valuable for stations with short records "
    "(<20 years) where AMS gives too few data points for reliable fitting.",
    font_size=14, color=DARK_GRAY)

add_para(tf, "", font_size=5)

# Two columns
left_items = [
    ("Threshold selection",
     "Drag a slider to set the threshold as a percentile (e.g. Q90) or type a "
     "value in m³/s directly. The chart updates live."),
    ("Independence filter",
     "A minimum separation gap (3, 7, 14, or 30 days) ensures only independent "
     "peaks are included — no double-counting of the same event."),
    ("GPD fitting",
     "The Generalised Pareto Distribution (GPD) is fitted by L-moments "
     "(closed-form, no iteration). Shape and scale parameters are reported."),
]
right_items = [
    ("Frequency curve",
     "Empirical exceedance plotted against fitted GPD line on a log return-period axis. "
     "Poisson process assumption links threshold-crossing rate to return periods."),
    ("Design table",
     "Quantiles for T = 2–500 years alongside the AMS results for direct comparison."),
    ("Warning system",
     "Automatic alert if fewer than 15 peaks are selected — advises increasing "
     "the record length or lowering the threshold."),
]

for i, (title, desc) in enumerate(left_items):
    add_para(tf, f"  ▸  {title}:", font_size=14, bold=True, color=NAVY, space_before=Pt(8))
    add_para(tf, f"      {desc}", font_size=12, color=DARK_GRAY, space_before=Pt(2))


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — PDF REPORT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
body_bg(s)
slide_header(s, "Printable PDF Report",
             "One-click export of a complete station analysis for client deliverables")

# Left column
add_rect(s, Inches(0.4), Inches(1.45), Inches(5.9), Inches(5.6), WHITE)
txb = s.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(5.55), Inches(5.2))
tf  = txb.text_frame; tf.word_wrap = True
add_para(tf, "🖨️  Print / Save PDF", font_size=18, bold=True, color=TEAL, space_before=Pt(0))
add_para(tf,
    "Every station has a dedicated report page that assembles all key "
    "analyses into a single printable document. Click the Print button "
    "to open the browser's PDF dialog and save.",
    font_size=14, color=DARK_GRAY)
add_para(tf, "", font_size=5)
add_para(tf, "Report includes:", font_size=14, bold=True, color=TEAL)
sections = [
    "Station metadata (ID, province, drainage area, record span)",
    "Full hydrograph",
    "Flow Duration Curve",
    "Annual Regime (mean + Q10–Q90 band)",
    "Annual peaks bar chart",
    "Flood frequency curve (best-fit distribution)",
    "Design flood quantile table",
    "Annual statistics table (year-by-year)",
    "Trend analysis summary paragraph",
    "Data attribution & disclaimer footer",
]
for sec in sections:
    add_para(tf, f"  • {sec}", font_size=12, color=DARK_GRAY, space_before=Pt(4))

# Right column
add_rect(s, Inches(6.55), Inches(1.45), Inches(6.4), Inches(5.6), NAVY)
txb = s.shapes.add_textbox(Inches(6.75), Inches(1.6), Inches(6.0), Inches(5.2))
tf  = txb.text_frame; tf.word_wrap = True
add_para(tf, "Report Design", font_size=18, bold=True, color=WHITE, space_before=Pt(0))
add_para(tf, "", font_size=5)
for item in [
    ("Print-optimised layout",
     "All interactive elements (buttons, nav, methodology panels) are automatically "
     "hidden in print mode via CSS @media print rules."),
    ("Static charts",
     "Plotly charts switch to a static, high-fidelity render for print — "
     "no toolbars or hover overlays."),
    ("Page breaks",
     "Each analysis section starts on a fresh page so the report "
     "flows cleanly across A4 or Letter paper."),
    ("URL to share",
     "The report URL (/stations/XXXXX/report) can be bookmarked "
     "or sent directly to colleagues."),
]:
    add_para(tf, f"▸  {item[0]}", font_size=14, bold=True, color=LIGHT_TEAL, space_before=Pt(10))
    add_para(tf, item[1], font_size=12, color=WHITE, space_before=Pt(3))


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — METHODOLOGY TRANSPARENCY
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
body_bg(s)
slide_header(s, "Methodology Transparency",
             "Every analysis tab includes an expandable explanation of the method used")

add_rect(s, Inches(0.4), Inches(1.45), Inches(12.5), Inches(5.6), WHITE)
txb = s.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(11.9), Inches(5.2))
tf  = txb.text_frame; tf.word_wrap = True

add_para(tf,
    "Each analysis tab has an 'ℹ About this analysis' accordion panel at the bottom of the page. "
    "Expanding it reveals a plain-English explanation of the method, the assumptions made, "
    "and peer-reviewed references — so any team member can understand exactly how a result was produced.",
    font_size=15, color=DARK_GRAY, space_before=Pt(0))

add_para(tf, "", font_size=5)
add_para(tf, "Covered in every tab:", font_size=15, bold=True, color=TEAL)

method_tabs = [
    ("Hydrograph",         "What daily mean discharge represents; ECCC data collection methods"),
    ("Flow Duration Curve","Exceedance probability calculation; interpretation of Q10/Q50/Q90"),
    ("Annual Regime",      "Climatological averaging; interpretation of seasonality"),
    ("Annual Statistics",  "How annual maxima, minima, and volumes are derived"),
    ("Low Flow",           "7Q10 derivation; relevance to environmental flows"),
    ("Trend Analysis",     "Mann-Kendall theory; Sen's slope; significance testing"),
    ("Flood Frequency",    "AMS approach; L-moments vs MLE; AIC model selection"),
    ("POT Analysis",       "GPD theory; Poisson process assumption; threshold selection guidance"),
    ("Baseflow",           "Lyne-Hollick and Eckhardt filter equations; BFI interpretation"),
    ("Goodness of Fit",    "K-S test; AIC/BIC; how to interpret scores"),
]
for tab, desc in method_tabs:
    add_para(tf, f"  • {tab}:  {desc}", font_size=12, color=DARK_GRAY, space_before=Pt(5))


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — TECHNOLOGY STACK
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
body_bg(s)
slide_header(s, "Technology Stack",
             "Modern, open-source, cloud-hosted — no licences required")

cards_data = [
    (Inches(0.4),  Inches(1.5),  Inches(3.9), Inches(2.4),
     "🖥️", "Frontend",
     ["Next.js 16 (React)", "TypeScript", "Tailwind CSS", "Plotly.js charts", "Leaflet maps"]),
    (Inches(4.55), Inches(1.5),  Inches(3.9), Inches(2.4),
     "⚙️", "Backend (FFA Engine)",
     ["Python 3.12", "FastAPI", "lmoments3", "NumPy / SciPy", "Docker on Render"]),
    (Inches(8.7),  Inches(1.5),  Inches(4.25), Inches(2.4),
     "🌐", "Data Source",
     ["ECCC OGC API (live)", "National HYDAT database", "No local DB needed", "Updated daily by WSC"]),
    (Inches(0.4),  Inches(4.1),  Inches(3.9), Inches(2.5),
     "☁️", "Hosting",
     ["Vercel (Next.js)", "Render (Python API)", "GitHub (source control)", "Free tier — $0/month"]),
    (Inches(4.55), Inches(4.1),  Inches(3.9), Inches(2.5),
     "🔒", "Architecture",
     ["Server-side API proxy", "No CORS exposure", "Environment-variable secrets", "Stateless — no user data stored"]),
    (Inches(8.7),  Inches(4.1),  Inches(4.25), Inches(2.5),
     "📐", "Analysis (client-side)",
     ["POT / GPD in TypeScript", "Baseflow filters in TypeScript", "Mann-Kendall in TypeScript", "FFA distributions via Python API"]),
]

for x, y, w, h, icon, heading, bullets in cards_data:
    add_rect(s, x, y, w, h, WHITE, line_rgb=TEAL, line_width=Pt(1.5))
    add_text(s, f"{icon}  {heading}",
             x+Inches(0.12), y+Inches(0.1), w-Inches(0.2), Inches(0.38),
             font_size=14, bold=True, color=TEAL)
    by = y + Inches(0.5)
    for b in bullets:
        add_text(s, f"  • {b}",
                 x+Inches(0.12), by, w-Inches(0.2), Inches(0.3),
                 font_size=11, color=DARK_GRAY)
        by += Inches(0.3)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — CLOSING / CTA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)

add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, 0, 0, Inches(0.18), SLIDE_H, TEAL)
add_rect(s, 0, Inches(6.3), SLIDE_W, Inches(1.2), TEAL)

add_text(s, "Try It Now",
         Inches(0.55), Inches(1.0), Inches(12), Inches(1.2),
         font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s,
    "The app is live — no login, no installation.",
    Inches(0.55), Inches(2.2), Inches(12), Inches(0.7),
    font_size=20, color=LIGHT_TEAL, align=PP_ALIGN.CENTER, italic=True)

# URL badge
add_rect(s, Inches(3.9), Inches(3.1), Inches(5.5), Inches(0.85), ORANGE)
add_text(s, "🌐  ws-cprep.vercel.app",
         Inches(4.0), Inches(3.15), Inches(5.3), Inches(0.75),
         font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s,
    "Search for any Canadian WSC station\nand all analyses run instantly in your browser.",
    Inches(1.5), Inches(4.25), Inches(10), Inches(1.2),
    font_size=17, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s,
    "Source code:  github.com/adilch/WSCprep",
    Inches(1.5), Inches(5.35), Inches(10), Inches(0.55),
    font_size=14, color=LIGHT_TEAL, align=PP_ALIGN.CENTER)

add_text(s,
    "Data: Environment and Climate Change Canada · National Hydrometric Program · HYDAT database",
    Inches(0.3), Inches(6.35), Inches(12.7), Inches(0.6),
    font_size=13, color=WHITE, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "WSCprep_Introduction.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
